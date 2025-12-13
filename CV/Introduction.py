from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_kids_intro():
    prs = Presentation()

    def add_slide(title_text, content_text, layout_index=1):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        
        # Title Styling
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.name = "Comic Sans MS"
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192) # Fun Blue
        title.text_frame.paragraphs[0].font.bold = True

        # Content Styling
        if layout_index == 1: # Standard Content
            content = slide.placeholders[1]
            content.text = content_text
            for paragraph in content.text_frame.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.size = Pt(28) # Big text for kids
        
        return slide

    # --- SLIDE 1: Intro / Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title Slide
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "HELLO EVERYONE! 👋"
    subtitle.text = "My Name is Max!\nLet's be friends!"
    
    # Custom Title Style
    title.text_frame.paragraphs[0].font.name = "Comic Sans MS"
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 102, 0) # Orange

    # --- SLIDE 2: Who Am I? ---
    add_slide(
        "Who Am I? 🤔",
        "• My name is Max!\n"
        "• I am 20 years old 🎂\n"
        "• I go to a big school called 'Asia University'\n"
        "• I learn how to talk to computers! 💻"
    )

    # --- SLIDE 3: Where I'm From ---
    add_slide(
        "Where Do I Live? 🌏",
        "• I live here in Taiwan now!\n"
        "• But I was born in INDONESIA 🇮🇩\n"
        "• It is a beautiful place with lots of islands.\n"
        "• Can you say 'Halo'? (That means Hello!)"
    )

    # --- SLIDE 4: My Hobbies ---
    add_slide(
        "What Do I Like To Do? 🎮",
        "1. I love playing Video Games 🕹️\n"
        "2. I love reading cool stories 📚\n"
        "3. I love eating yummy food! 🍜\n\n"
        "What is YOUR favorite food?"
    )

    # --- SLIDE 5: Quiz Time ---
    add_slide(
        "Pop Quiz! ❓",
        "Let's see if you were listening!\n\n"
        "Q1: How old am I?\n"
        "Q2: What country am I from?\n"
        "Q3: Do I like video games?"
    )

    # --- SLIDE 6: Questions ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Any Questions? 🙋‍♂️"
    subtitle.text = "Ask me anything!"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 153, 0) # Green

    # Save
    filename = "Max_Intro_For_Kids.pptx"
    prs.save(filename)
    print(f"✅ Success! Created '{filename}'")

if __name__ == "__main__":
    create_kids_intro()